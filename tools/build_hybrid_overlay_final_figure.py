from base64 import b64encode
from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


W, H = 1672, 941
SRC = Path(
    "intermediate_artifacts/paper_framework_figure_studio/"
    "third_chapter_method_framework/outputs/S5-candidate-images/candidate-03.png"
)
OUT_DIR = Path(
    "intermediate_artifacts/paper_framework_figure_studio/"
    "third_chapter_method_framework/outputs/editable-final-figure-hybrid-v2"
)
SVG_PATH = OUT_DIR / "hybrid_overlay_editable_final_framework.svg"
PPTX_PATH = OUT_DIR / "hybrid_overlay_editable_final_framework.pptx"
MANIFEST_PATH = OUT_DIR / "hybrid_overlay_conversion_manifest.md"

BLUE = "#1557A8"
BLUE_DARK = "#0B3F86"
INK = "#1D2733"
RED = "#D42E20"
WHITE = "#FFFFFF"


def ensure_dir():
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def data_uri(path):
    return "data:image/png;base64," + b64encode(path.read_bytes()).decode("ascii")


def svg_rect(x, y, w, h, fill=WHITE, stroke="none", sw=0, rx=0, opacity=1):
    return (
        f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="{rx}" '
        f'fill="{fill}" fill-opacity="{opacity}" stroke="{stroke}" stroke-width="{sw}"/>'
    )


def svg_text(x, y, text, size=16, weight="400", fill=INK, anchor="middle"):
    style = f"font-family:Arial, Helvetica, sans-serif;font-size:{size}px;font-weight:{weight};fill:{fill};"
    return f'<text x="{x}" y="{y}" text-anchor="{anchor}" style="{style}">{escape(text)}</text>'


def svg_multiline(x, y, lines, size=15, gap=18, weight="400", fill=INK, anchor="middle"):
    return "\n".join(
        svg_text(x, y + i * gap, line, size=size, weight=weight, fill=fill, anchor=anchor)
        for i, line in enumerate(lines)
    )


def svg_line(x1, y1, x2, y2, stroke=INK, sw=2, dash=None, marker=False):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    marker_attr = ' marker-end="url(#arrow)"' if marker else ""
    return f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="{stroke}" stroke-width="{sw}"{dash_attr}{marker_attr}/>'


def text_specs():
    # mask_x, mask_y, mask_w, mask_h, text_x, text_y, lines, font_size, gap, weight, color
    specs = [
        (62, 125, 110, 32, 107, 148, ["Input Data"], 20, 22, "700", INK),
        (290, 120, 92, 58, 332, 142, ["Planning-level", "Risk Fields"], 20, 23, "700", INK),
        (496, 118, 120, 78, 555, 139, ["DEM-driven", "Adaptive Flight", "Corridor"], 17, 21, "700", INK),
        (775, 120, 360, 34, 962, 143, ["Terrain-aware Three-layer Airway Network"], 21, 22, "700", INK),
        (1327, 118, 92, 70, 1370, 136, ["Event-driven", "Incremental", "Update"], 17, 20, "700", INK),
        (1535, 120, 78, 56, 1574, 139, ["Path", "Post-processing"], 18, 22, "700", INK),
        (134, 238, 52, 26, 160, 256, ["DEM"], 16, 18, "700", INK),
        (128, 313, 68, 48, 160, 332, ["OSM", "features"], 15, 18, "400", INK),
        (110, 431, 95, 48, 158, 450, ["Communication", "source (LOS)"], 13.5, 17, "400", INK),
        (107, 557, 82, 30, 148, 578, ["Terminals"], 15, 18, "400", INK),
        (105, 635, 82, 47, 146, 654, ["UAV", "parameters"], 15, 18, "400", INK),
        (333, 238, 82, 52, 374, 257, ["Terrain", "representation"], 15, 18, "400", INK),
        (328, 388, 90, 55, 372, 407, ["Human exposure", "risk proxy"], 14.5, 18, "400", INK),
        (326, 539, 92, 60, 373, 558, ["LOS communication", "risk proxy"], 13.5, 17, "400", INK),
        (463, 192, 120, 24, 523, 210, ["Floor boundary"], 14, 18, "400", INK),
        (468, 315, 128, 26, 531, 334, ["Ceiling boundary"], 14, 18, "400", INK),
        (462, 360, 170, 30, 547, 381, ["Three mid-surfaces"], 16, 18, "700", INK),
        (560, 414, 92, 24, 606, 432, ["Ceiling mid-surface"], 12.5, 16, "400", INK),
        (560, 444, 92, 24, 606, 462, ["Middle mid-surface"], 12.5, 16, "400", INK),
        (560, 474, 92, 24, 606, 492, ["Floor mid-surface"], 12.5, 16, "400", INK),
        (462, 520, 165, 26, 545, 539, ["Flyable mask (corridor)"], 15, 18, "700", INK),
        (912, 234, 95, 52, 960, 253, ["Endpoint access", "layer"], 14, 16, "400", BLUE),
        (914, 337, 98, 52, 964, 356, ["Regional branch", "layer"], 14, 16, "400", "#5E8E3E"),
        (914, 450, 82, 48, 955, 467, ["Backbone", "layer"], 14, 16, "400", RED),
        (1026, 170, 198, 24, 1125, 188, ["Corridor to graph mechanism"], 14.5, 17, "700", INK),
        (1047, 204, 118, 28, 1105, 224, ["Mid-surfaces input"], 12.5, 15, "400", INK),
        (1043, 266, 106, 42, 1096, 284, ["Terrain-driven", "node sampling"], 12.5, 15, "400", INK),
        (1043, 328, 106, 42, 1096, 346, ["Candidate edges", "(intra-layer)"], 12.5, 15, "400", INK),
        (1043, 390, 106, 42, 1096, 408, ["Safety check", "(clearance & risk)"], 12.5, 15, "400", INK),
        (1043, 452, 112, 42, 1098, 470, ["Inter-layer links", "(vertical connections)"], 12.0, 15, "400", INK),
        (1043, 514, 100, 42, 1094, 532, ["Cost assignment", "c(e)"], 12.5, 15, "400", INK),
        (1043, 576, 100, 42, 1094, 594, ["Graph output", "G = (V, E)"], 12.5, 15, "400", INK),
        (1363, 145, 72, 48, 1398, 164, ["Regional", "event"], 14.5, 18, "400", INK),
        (1360, 258, 90, 48, 1405, 277, ["Local affected", "edges"], 14, 18, "400", RED),
        (1346, 365, 104, 54, 1398, 384, ["LPA* state reuse", "(g, rhs)"], 13.5, 17, "400", INK),
        (1353, 590, 92, 28, 1398, 609, ["Updated path"], 14, 18, "400", INK),
        (1530, 250, 105, 24, 1584, 268, ["LOS pruning"], 14, 18, "400", INK),
        (1530, 397, 100, 42, 1580, 415, ["B-spline", "smoothing"], 13.5, 17, "400", INK),
        (1578, 478, 74, 50, 1614, 496, ["Safety recheck", "(clearance & risk)"], 12.5, 16, "400", INK),
        (1520, 620, 122, 62, 1582, 637, ["Continuous", "trajectory after", "path update"], 12.5, 15, "400", INK),
        (344, 720, 245, 32, 466, 742, ["Adaptive corridor generation"], 17, 18, "400", INK),
        (898, 720, 260, 32, 1028, 742, ["Structured graph compression"], 17, 18, "400", INK),
        (1362, 720, 250, 32, 1488, 742, ["Local update plus post-processing"], 17, 18, "400", INK),
        (235, 816, 110, 38, 290, 842, ["Outcome"], 22, 18, "700", INK),
        (380, 811, 1080, 45, 920, 842, ["Mountainous UAV replanning is transformed into a compact, safety-constrained, incrementally updateable graph search problem."], 20, 18, "400", INK),
    ]
    return specs


def build_svg():
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">',
        "<defs>",
        '<marker id="arrow" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="8" markerHeight="8" orient="auto-start-reverse">',
        f'<path d="M 0 0 L 10 5 L 0 10 z" fill="{INK}"/>',
        "</marker>",
        "</defs>",
        f'<image x="0" y="0" width="{W}" height="{H}" href="{data_uri(SRC)}"/>',
    ]
    for mx, my, mw, mh, *_ in text_specs():
        parts.append(svg_rect(mx, my, mw, mh, fill=WHITE, opacity=0.98))
    for _, _, _, _, tx, ty, lines, size, gap, weight, color in text_specs():
        parts.append(svg_multiline(tx, ty, lines, size=size, gap=gap, weight=weight, fill=color))
    # 加少量可编辑主流程箭头，覆盖原箭头，便于后续调粗细和位置。
    for x1, x2 in [(200, 236), (427, 452), (657, 674), (1248, 1273), (1467, 1490)]:
        parts.append(svg_line(x1, 384, x2, 384, stroke=INK, sw=5, marker=True))
    parts.append("</svg>")
    SVG_PATH.write_text("\n".join(parts), encoding="utf-8")


def pptx_rgb(hex_color):
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Canvas:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.sx = 13.333 / W
        self.sy = 7.5 / H

    def ix(self, x):
        return Inches(x * self.sx)

    def iy(self, y):
        return Inches(y * self.sy)

    def rect(self, x, y, w, h):
        shp = self.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = pptx_rgb(WHITE)
        shp.line.color.rgb = pptx_rgb(WHITE)
        shp.line.width = Pt(0)
        return shp

    def text(self, x, y, w, h, text, size=10, weight=False, color=INK):
        box = self.slide.shapes.add_textbox(self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        box.text_frame.clear()
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = weight
        run.font.color.rgb = pptx_rgb(color)
        return box

    def line(self, x1, y1, x2, y2):
        conn = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, self.ix(x1), self.iy(y1), self.ix(x2), self.iy(y2))
        conn.line.color.rgb = pptx_rgb(INK)
        conn.line.width = Pt(2.2)
        return conn


def build_pptx():
    c = Canvas()
    c.slide.shapes.add_picture(str(SRC), c.ix(0), c.iy(0), c.ix(W), c.iy(H))
    for mx, my, mw, mh, *_ in text_specs():
        c.rect(mx, my, mw, mh)
    for mx, my, mw, mh, tx, ty, lines, size, gap, weight, color in text_specs():
        # 使用 mask 框作为文本框范围，避免文本错位；字号从 SVG 像素换算到 PPT 点。
        c.text(mx, my, mw, mh, "\n".join(lines), size=max(5, size * 0.52), weight=weight == "700", color=color)
    for x1, x2 in [(200, 236), (427, 452), (657, 674), (1248, 1273), (1467, 1490)]:
        c.line(x1, 384, x2, 384)
    c.prs.save(PPTX_PATH)


def write_manifest():
    MANIFEST_PATH.write_text(
        "\n".join(
            [
                "# 混合可编辑化 v2 说明",
                "",
                f"源图：`{SRC}`",
                f"SVG：`{SVG_PATH.name}`",
                f"PPTX：`{PPTX_PATH.name}`",
                "",
                "## 修复点",
                "",
                "上一版采用多个局部贴片，部分裁剪框过紧，造成图标、网络或小曲线只显示一半。本版改为完整源图保真底层，不再裁切复杂绘制区域，因此不会出现复杂对象被截断的问题。",
                "",
                "## 可编辑层",
                "",
                "本版在完整源图上对主要文字区域放置白色遮罩，并覆盖为可编辑文本；主流程箭头也额外覆盖为可编辑线条。复杂图形、山地、风险面、三层网络、事件图标和后处理图标保持最终 PNG 原样。",
                "",
                "## 取舍",
                "",
                "复杂图形不再作为单独小贴片裁切，因而保真度最高；若需要单独移动某个复杂图标，则仍建议使用上一版贴片思路，但需要逐个手调裁剪框。本版优先解决截断和半显示问题。",
                "",
                "## 语义保持",
                "",
                "第 4 步仍突出三层航路网络，第 5 步红色局部边仍表示区域事件影响的受影响边集合，第 6 步仍为路径后处理，图内阶段编号保持 1 至 6。",
                "",
            ]
        ),
        encoding="utf-8",
    )


def main():
    ensure_dir()
    build_svg()
    build_pptx()
    write_manifest()
    print(SVG_PATH)
    print(PPTX_PATH)
    print(MANIFEST_PATH)


if __name__ == "__main__":
    main()
