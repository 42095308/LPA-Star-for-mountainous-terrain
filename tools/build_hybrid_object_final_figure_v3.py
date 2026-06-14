from base64 import b64encode
from pathlib import Path
from xml.sax.saxutils import escape

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


W, H = 1672, 941
SRC = Path(
    "intermediate_artifacts/paper_framework_figure_studio/"
    "third_chapter_method_framework/outputs/S5-candidate-images/candidate-03.png"
)
OUT_DIR = Path(
    "intermediate_artifacts/paper_framework_figure_studio/"
    "third_chapter_method_framework/outputs/editable-final-figure-hybrid-v4-crop-fix"
)
ASSET_DIR = OUT_DIR / "assets"
SVG_PATH = OUT_DIR / "hybrid_object_editable_final_framework_v4.svg"
PPTX_PATH = OUT_DIR / "hybrid_object_editable_final_framework_v4.pptx"
MANIFEST_PATH = OUT_DIR / "hybrid_object_conversion_manifest_v4.md"
CONTACT_PATH = OUT_DIR / "asset_contact_sheet.png"

BLUE = "#1557A8"
BLUE_DARK = "#0B3F86"
GREEN = "#5E8E3E"
RED = "#D42E20"
INK = "#1D2733"
GRID = "#D9E2EF"
WHITE = "#FFFFFF"


# 精裁原则：保留完整复杂图形，尽量不带文字；无法完全分离的对象采用更宽边界，后续用白色遮罩覆盖原文字并重建为可编辑文字。
CROPS = [
    ("stage1_dem", 25, 183, 116, 78),
    ("stage1_osm", 32, 278, 92, 82),
    ("stage1_los", 22, 374, 76, 82),
    ("stage1_terminals", 22, 480, 150, 44),
    ("stage1_uav", 24, 556, 58, 78),
    ("stage2_terrain", 238, 199, 116, 82),
    ("stage2_exposure", 238, 321, 88, 92),
    ("stage2_los", 238, 448, 88, 105),
    ("stage3_floor", 468, 214, 162, 70),
    ("stage3_ceiling", 468, 313, 162, 72),
    ("stage3_mids", 466, 407, 96, 96),
    ("stage3_mask", 462, 550, 170, 94),
    ("stage4_network", 682, 168, 325, 465),
    ("stage4_mech_1", 1144, 210, 78, 50),
    ("stage4_mech_2", 1132, 277, 94, 58),
    ("stage4_mech_3", 1132, 346, 94, 54),
    ("stage4_mech_4", 1144, 410, 72, 72),
    ("stage4_mech_5", 1144, 483, 80, 56),
    ("stage4_mech_6", 1150, 548, 72, 42),
    ("stage4_mech_7", 1130, 590, 96, 72),
    ("stage5_event", 1293, 202, 66, 90),
    ("stage5_edges", 1281, 318, 82, 96),
    ("stage5_state", 1278, 450, 78, 74),
    ("stage5_path", 1284, 570, 92, 58),
    ("stage6_los", 1498, 205, 138, 40),
    ("stage6_bspline", 1498, 322, 142, 52),
    ("stage6_shield", 1498, 454, 66, 76),
    ("stage6_traj", 1496, 546, 142, 52),
    ("callout_corridor", 224, 708, 118, 62),
    ("callout_graph", 748, 708, 90, 66),
    ("callout_update", 1252, 708, 106, 66),
    ("outcome_target", 118, 798, 108, 82),
]


CLEAN_MASKS = {
    "stage2_terrain": [(96, 50, 36, 32)],
    "stage3_mids": [(78, 0, 28, 96)],
    "stage4_network": [(230, 56, 94, 54), (230, 162, 96, 56), (231, 272, 88, 54), (3, 418, 305, 42)],
    "stage6_shield": [(42, 0, 36, 76)],
    "stage6_bspline": [(45, 44, 80, 14)],
    "stage6_traj": [(50, 44, 86, 14)],
}


def ensure_dirs():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def crop_assets():
    im = Image.open(SRC).convert("RGBA")
    assets = {}
    for name, x, y, w, h in CROPS:
        crop = im.crop((x, y, x + w, y + h))
        if name in CLEAN_MASKS:
            draw = ImageDraw.Draw(crop)
            for mx, my, mw, mh in CLEAN_MASKS[name]:
                draw.rectangle((mx, my, mx + mw, my + mh), fill=(255, 255, 255, 255))
        path = ASSET_DIR / f"{name}.png"
        crop.save(path)
        assets[name] = {"path": path, "x": x, "y": y, "w": w, "h": h}
    return assets


def image_data_uri(path):
    return "data:image/png;base64," + b64encode(path.read_bytes()).decode("ascii")


def svg_rect(x, y, w, h, rx=0, fill=WHITE, stroke=BLUE_DARK, sw=2, dash=None):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    return f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="{rx}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}"{dash_attr}/>'


def svg_line(x1, y1, x2, y2, stroke=INK, sw=2, dash=None, marker=False):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    marker_attr = ' marker-end="url(#arrow)"' if marker else ""
    return f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="{stroke}" stroke-width="{sw}"{dash_attr}{marker_attr}/>'


def svg_circle(cx, cy, r, fill=BLUE, stroke="none", sw=1):
    return f'<circle cx="{cx}" cy="{cy}" r="{r}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}"/>'


def svg_text(x, y, text, size=16, weight="400", fill=INK, anchor="middle"):
    style = f"font-family:Arial, Helvetica, sans-serif;font-size:{size}px;font-weight:{weight};fill:{fill};"
    return f'<text x="{x}" y="{y}" text-anchor="{anchor}" style="{style}">{escape(text)}</text>'


def svg_multiline(x, y, lines, size=15, gap=18, weight="400", fill=INK, anchor="middle"):
    return "\n".join(
        svg_text(x, y + i * gap, line, size=size, weight=weight, fill=fill, anchor=anchor)
        for i, line in enumerate(lines)
    )


def svg_image(asset):
    return (
        f'<image x="{asset["x"]}" y="{asset["y"]}" width="{asset["w"]}" height="{asset["h"]}" '
        f'href="{image_data_uri(asset["path"])}"/>'
    )


def add_header(parts, x, y, w, number, title_lines):
    parts.append(svg_circle(x + 20, y + 6, 20, BLUE))
    parts.append(svg_text(x + 20, y + 14, str(number), size=25, weight="700", fill=WHITE))
    parts.append(svg_multiline(x + w / 2, y + 42, title_lines, size=20, gap=24, weight="700"))


def add_editable_stage4_legend(parts):
    # 覆盖 stage4_network 内部的原始层名和底部 legend，重建为可编辑对象。
    masks = [
        (892, 214, 130, 70),
        (892, 320, 130, 72),
        (892, 430, 130, 72),
        (685, 586, 300, 40),
    ]
    for m in masks:
        parts.append(svg_rect(*m, rx=0, fill=WHITE, stroke=WHITE, sw=0))
    parts.append(svg_multiline(960, 250, ["Endpoint access", "layer"], size=14, gap=16, fill=BLUE))
    parts.append(svg_multiline(964, 356, ["Regional branch", "layer"], size=14, gap=16, fill=GREEN))
    parts.append(svg_multiline(955, 466, ["Backbone", "layer"], size=14, gap=16, fill=RED))
    parts.append(svg_circle(698, 605, 6, BLUE))
    parts.append(svg_text(708, 609, "Endpoint node", size=12, fill=INK, anchor="start"))
    parts.append(svg_circle(802, 605, 6, GREEN))
    parts.append(svg_text(812, 609, "Branch node", size=12, fill=INK, anchor="start"))
    parts.append(svg_circle(902, 605, 6, RED))
    parts.append(svg_text(912, 609, "Backbone node", size=12, fill=INK, anchor="start"))
    parts.append(svg_line(785, 623, 820, 623, stroke=INK, sw=1.2, dash="6 4"))
    parts.append(svg_text(828, 627, "Inter-layer link", size=12, fill=INK, anchor="start"))


def build_svg(assets):
    cards = [
        (15, 104, 185, 555),
        (236, 104, 191, 555),
        (452, 104, 205, 555),
        (674, 104, 574, 555),
        (1273, 104, 194, 555),
        (1490, 104, 168, 555),
    ]
    titles = [
        ["Input Data"],
        ["Planning-level", "Risk Fields"],
        ["DEM-driven", "Adaptive Flight", "Corridor"],
        ["Terrain-aware Three-layer Airway Network"],
        ["Event-driven", "Incremental Update"],
        ["Path", "Post-processing"],
    ]

    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">',
        "<defs>",
        '<marker id="arrow" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="8" markerHeight="8" orient="auto-start-reverse">',
        f'<path d="M 0 0 L 10 5 L 0 10 z" fill="{INK}"/>',
        "</marker>",
        "</defs>",
        svg_rect(0, 0, W, H, fill=WHITE, stroke=WHITE),
    ]
    for i, (x, y, w, h) in enumerate(cards):
        parts.append(svg_rect(x, y, w, h, rx=12, fill=WHITE, stroke=BLUE_DARK, sw=2))
        add_header(parts, x, y, w, i + 1, titles[i])
    for x1, x2 in [(200, 236), (427, 452), (657, 674), (1248, 1273), (1467, 1490)]:
        parts.append(svg_line(x1, 384, x2, 384, stroke=INK, sw=5, marker=True))

    # 可编辑分隔线
    for x, ys in [(15, [185, 300, 430, 550]), (1273, [290, 430, 558]), (1490, [290, 430, 558])]:
        ww = 185 if x == 15 else 194 if x == 1273 else 168
        for yy in ys:
            parts.append(svg_line(x + 12, yy, x + ww - 12, yy, stroke=GRID, sw=1.3))

    # 放置所有复杂 raster 对象
    order = [
        "stage1_dem", "stage1_osm", "stage1_los", "stage1_terminals", "stage1_uav",
        "stage2_terrain", "stage2_exposure", "stage2_los",
        "stage3_floor", "stage3_ceiling", "stage3_mids", "stage3_mask",
        "stage4_network",
        "stage4_mech_1", "stage4_mech_2", "stage4_mech_3", "stage4_mech_4", "stage4_mech_5", "stage4_mech_6", "stage4_mech_7",
        "stage5_event", "stage5_edges", "stage5_state", "stage5_path",
        "stage6_los", "stage6_bspline", "stage6_shield", "stage6_traj",
        "callout_corridor", "callout_graph", "callout_update", "outcome_target",
    ]
    for name in order:
        parts.append(svg_image(assets[name]))

    # Stage 4 机制面板框与每一步文字可编辑
    parts.append(svg_rect(1014, 168, 215, 480, rx=8, fill="none", stroke=BLUE_DARK, sw=1.8))
    parts.append(svg_text(1122, 188, "Corridor to graph mechanism", size=14.5, weight="700"))
    mech_labels = [
        (1105, 224, ["Mid-surfaces input"]),
        (1096, 284, ["Terrain-driven", "node sampling"]),
        (1096, 346, ["Candidate edges", "(intra-layer)"]),
        (1096, 408, ["Safety check", "(clearance & risk)"]),
        (1098, 470, ["Inter-layer links", "(vertical connections)"]),
        (1094, 532, ["Cost assignment", "c(e)"]),
        (1094, 594, ["Graph output", "G = (V, E)"]),
    ]
    for i, (tx, ty, lines) in enumerate(mech_labels):
        by = 206 + i * 62
        parts.append(svg_rect(1024, by - 2, 195, 55, rx=6, fill="none", stroke="#9CBDE4", sw=1))
        parts.append(svg_circle(1037, by + 16, 10, BLUE))
        parts.append(svg_text(1037, by + 21, str(i + 1), size=12, weight="700", fill=WHITE))
        parts.append(svg_multiline(tx, ty, lines, size=12.5, gap=15))
        if i < 6:
            parts.append(svg_line(1121, by + 53, 1121, by + 62, stroke=INK, sw=1.1, marker=True))

    add_editable_stage4_legend(parts)

    # 可编辑文字
    text_blocks = [
        (160, 256, ["DEM"], 16, 18, "700", INK),
        (160, 332, ["OSM", "features"], 15, 18, "400", INK),
        (158, 450, ["Communication", "source (LOS)"], 13.5, 17, "400", INK),
        (148, 578, ["Terminals"], 15, 18, "400", INK),
        (146, 654, ["UAV", "parameters"], 15, 18, "400", INK),
        (374, 257, ["Terrain", "representation"], 15, 18, "400", INK),
        (372, 407, ["Human exposure", "risk proxy"], 14.5, 18, "400", INK),
        (373, 558, ["LOS communication", "risk proxy"], 13.5, 17, "400", INK),
        (523, 210, ["Floor boundary"], 14, 18, "400", INK),
        (531, 334, ["Ceiling boundary"], 14, 18, "400", INK),
        (547, 381, ["Three mid-surfaces"], 16, 18, "700", INK),
        (606, 432, ["Ceiling mid-surface"], 12.5, 16, "400", INK),
        (606, 462, ["Middle mid-surface"], 12.5, 16, "400", INK),
        (606, 492, ["Floor mid-surface"], 12.5, 16, "400", INK),
        (545, 539, ["Flyable mask (corridor)"], 15, 18, "700", INK),
        (1398, 226, ["Regional event"], 13.5, 18, "400", INK),
        (1405, 277, ["Local affected", "edges"], 14, 18, "400", RED),
        (1398, 384, ["LPA* state reuse", "(g, rhs)"], 13.5, 17, "400", INK),
        (1398, 638, ["Updated path"], 12.5, 18, "400", INK),
        (1584, 268, ["LOS pruning"], 14, 18, "400", INK),
        (1580, 415, ["B-spline", "smoothing"], 13.5, 17, "400", INK),
        (1614, 496, ["Safety recheck", "(clearance & risk)"], 12.5, 16, "400", INK),
        (1582, 621, ["Continuous trajectory", "after path update"], 11.5, 14, "400", INK),
        (466, 742, ["Adaptive corridor generation"], 17, 18, "400", INK),
        (1028, 742, ["Structured graph compression"], 17, 18, "400", INK),
        (1488, 742, ["Local update plus post-processing"], 17, 18, "400", INK),
        (290, 842, ["Outcome"], 22, 18, "700", INK),
        (920, 842, ["Mountainous UAV replanning is transformed into a compact, safety-constrained, incrementally updatable graph search problem."], 20, 18, "400", INK),
    ]
    for tx, ty, lines, size, gap, weight, color in text_blocks:
        parts.append(svg_multiline(tx, ty, lines, size=size, gap=gap, weight=weight, fill=color))

    # 底部框
    for cx, cy, cw, ch in [(207, 695, 409, 80), (710, 695, 452, 80), (1245, 695, 391, 80)]:
        parts.append(svg_rect(cx, cy, cw, ch, rx=9, fill="none", stroke=BLUE_DARK, sw=2, dash="8 5"))
    parts.append(svg_rect(91, 792, 1460, 86, rx=11, fill="none", stroke=BLUE_DARK, sw=2))
    parts.append(svg_line(330, 806, 330, 864, stroke=BLUE_DARK, sw=2))
    parts.append("</svg>")
    SVG_PATH.write_text("\n".join(parts), encoding="utf-8")


def pptx_rgb(hex_color):
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Canvas:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.sx = 13.333 / W
        self.sy = 7.5 / H

    def ix(self, x):
        return Inches(x * self.sx)

    def iy(self, y):
        return Inches(y * self.sy)

    def rect(self, x, y, w, h, fill=WHITE, stroke=BLUE_DARK, sw=1.0, radius=True):
        typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shp = self.slide.shapes.add_shape(typ, self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = pptx_rgb(fill)
        shp.line.color.rgb = pptx_rgb(stroke)
        shp.line.width = Pt(sw)
        return shp

    def oval(self, x, y, w, h, fill=BLUE):
        shp = self.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = pptx_rgb(fill)
        shp.line.color.rgb = pptx_rgb(fill)
        return shp

    def image(self, asset):
        self.slide.shapes.add_picture(str(asset["path"]), self.ix(asset["x"]), self.iy(asset["y"]), self.ix(asset["w"]), self.iy(asset["h"]))

    def text(self, x, y, w, h, text, size=10, weight=False, color=INK):
        box = self.slide.shapes.add_textbox(self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        box.text_frame.clear()
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = weight
        run.font.color.rgb = pptx_rgb(color)
        return box

    def line(self, x1, y1, x2, y2, color=INK, sw=1.2):
        conn = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, self.ix(x1), self.iy(y1), self.ix(x2), self.iy(y2))
        conn.line.color.rgb = pptx_rgb(color)
        conn.line.width = Pt(sw)
        return conn


def build_pptx(assets):
    c = Canvas()
    c.rect(0, 0, W, H, fill=WHITE, stroke=WHITE, radius=False)
    cards = [
        (15, 104, 185, 555, "Input Data"),
        (236, 104, 191, 555, "Planning-level\nRisk Fields"),
        (452, 104, 205, 555, "DEM-driven\nAdaptive Flight\nCorridor"),
        (674, 104, 574, 555, "Terrain-aware Three-layer Airway Network"),
        (1273, 104, 194, 555, "Event-driven\nIncremental Update"),
        (1490, 104, 168, 555, "Path\nPost-processing"),
    ]
    for i, (x, y, w, h, title) in enumerate(cards):
        c.rect(x, y, w, h, fill=WHITE, stroke=BLUE_DARK, sw=1.0)
        c.oval(x, y - 18, 42, 42, fill=BLUE)
        c.text(x + 4, y - 15, 34, 35, str(i + 1), size=15, weight=True, color=WHITE)
        c.text(x + 10, y + 18, w - 20, 70, title, size=11.5, weight=True)
    for x1, x2 in [(200, 236), (427, 452), (657, 674), (1248, 1273), (1467, 1490)]:
        c.line(x1, 384, x2, 384, color=INK, sw=2.2)
    for name in assets:
        c.image(assets[name])
    # 文字用简化位置，主要供 PowerPoint 中后续人工编辑。
    text_items = [
        (140, 238, 70, 30, "DEM", 8, True, INK),
        (128, 313, 68, 45, "OSM\nfeatures", 7.8, False, INK),
        (104, 431, 100, 50, "Communication\nsource (LOS)", 6.8, False, INK),
        (110, 557, 80, 30, "Terminals", 8, False, INK),
        (106, 635, 80, 45, "UAV\nparameters", 8, False, INK),
        (330, 238, 88, 52, "Terrain\nrepresentation", 7.8, False, INK),
        (326, 388, 92, 56, "Human exposure\nrisk proxy", 7.5, False, INK),
        (326, 539, 92, 60, "LOS communication\nrisk proxy", 7.0, False, INK),
        (463, 192, 122, 25, "Floor boundary", 7.5, False, INK),
        (468, 315, 128, 25, "Ceiling boundary", 7.5, False, INK),
        (462, 360, 170, 30, "Three mid-surfaces", 8.4, True, INK),
        (452, 520, 185, 28, "Flyable mask (corridor)", 8.0, True, INK),
        (912, 224, 93, 54, "Endpoint access\nlayer", 7.2, False, BLUE),
        (912, 330, 96, 56, "Regional branch\nlayer", 7.2, False, GREEN),
        (913, 440, 84, 54, "Backbone\nlayer", 7.2, False, RED),
        (1024, 170, 195, 24, "Corridor to graph mechanism", 7.2, True, INK),
        (1363, 204, 72, 30, "Regional event", 7.0, False, INK),
        (1360, 258, 90, 48, "Local affected\nedges", 7.5, False, RED),
        (1346, 365, 104, 54, "LPA* state reuse\n(g, rhs)", 7.0, False, INK),
        (1353, 620, 92, 28, "Updated path", 6.6, False, INK),
        (1530, 250, 105, 24, "LOS pruning", 7.5, False, INK),
        (1530, 397, 100, 42, "B-spline\nsmoothing", 7.2, False, INK),
        (1578, 478, 74, 50, "Safety recheck\n(clearance & risk)", 6.2, False, INK),
        (1520, 607, 122, 50, "Continuous trajectory\nafter path update", 6.0, False, INK),
        (344, 720, 245, 32, "Adaptive corridor generation", 9, False, INK),
        (898, 720, 260, 32, "Structured graph compression", 9, False, INK),
        (1362, 720, 250, 32, "Local update plus post-processing", 9, False, INK),
        (235, 816, 110, 38, "Outcome", 12, True, INK),
        (380, 811, 1080, 45, "Mountainous UAV replanning is transformed into a compact, safety-constrained, incrementally updatable graph search problem.", 10.3, False, INK),
    ]
    # 覆盖 stage4 内部文字原图后再加 editable text，避免重复。
    for x, y, w, h, *_ in text_items:
        if x >= 912 and y < 500:
            c.rect(x, y, w, h, fill=WHITE, stroke=WHITE, sw=0, radius=False)
    for x, y, w, h, text, size, bold, color in text_items:
        c.text(x, y, w, h, text, size=size, weight=bold, color=color)
    c.prs.save(PPTX_PATH)


def make_contact_sheet(assets):
    files = [assets[name]["path"] for name, *_ in CROPS]
    tw, th = 260, 170
    cols = 4
    rows = (len(files) + cols - 1) // cols
    sheet = Image.new("RGB", (tw * cols, th * rows), WHITE)
    draw = ImageDraw.Draw(sheet)
    for i, path in enumerate(files):
        im = Image.open(path).convert("RGB")
        im.thumbnail((tw - 20, th - 42))
        x = (i % cols) * tw + 10
        y = (i // cols) * th + 8
        sheet.paste(im, (x, y))
        draw.text((x, y + th - 30), path.stem, fill=(0, 0, 0))
    sheet.save(CONTACT_PATH)


def write_manifest(assets):
    lines = [
        "# 对象级混合可编辑化 v3 说明",
        "",
        f"源图：`{SRC}`",
        f"SVG：`{SVG_PATH.name}`",
        f"PPTX：`{PPTX_PATH.name}`",
        f"贴片总览：`{CONTACT_PATH.name}`",
        "",
        "## 转换策略",
        "",
        "本版不再使用整张 PNG 作为底图，而是把不能稳定矢量化的复杂绘制部分分别裁切为完整 raster 对象。可稳定转换的阶段框、编号、主箭头、分隔线、标题和主要文字均重建为可编辑对象。",
        "",
        "## 对象级 raster 贴片",
        "",
    ]
    for name, asset in assets.items():
        lines.append(f"- `{name}.png`，x={asset['x']}，y={asset['y']}，w={asset['w']}，h={asset['h']}")
    lines.extend(
        [
            "",
            "## 相比 v1 的修复",
            "",
            "所有裁剪框均增加安全边距，三层网络、机制链小图、事件图标和路径后处理图标不再采用过窄裁剪。第 4 步网络主体作为一个完整对象保留，同时覆盖并重建层名和底部 legend，使文字仍可编辑。",
            "",
            "## 相比 v2 的区别",
            "",
            "v2 使用完整源图作为底层，保真但复杂元素不能单独移动。本版改为多个独立复杂对象，可以单独移动、替换、裁剪或删除，同时保留文字和框线的可编辑性。",
            "",
        ]
    )
    MANIFEST_PATH.write_text("\n".join(lines), encoding="utf-8")


def main():
    ensure_dirs()
    assets = crop_assets()
    build_svg(assets)
    build_pptx(assets)
    make_contact_sheet(assets)
    write_manifest(assets)
    print(SVG_PATH)
    print(PPTX_PATH)
    print(MANIFEST_PATH)
    print(CONTACT_PATH)
    print(f"assets={len(assets)}")


if __name__ == "__main__":
    main()
