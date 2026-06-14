from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


W, H = 1672, 941
OUT_DIR = Path(
    "intermediate_artifacts/paper_framework_figure_studio/"
    "third_chapter_method_framework/outputs/editable-final-figure"
)
SVG_PATH = OUT_DIR / "editable_final_framework.svg"
PPTX_PATH = OUT_DIR / "editable_final_framework.pptx"
MANIFEST_PATH = OUT_DIR / "editable_conversion_manifest.md"

BLUE = "#1557A8"
BLUE_DARK = "#0B3F86"
GREEN = "#5E8E3E"
RED = "#D42E20"
ORANGE = "#E5A33A"
INK = "#1D2733"
MUTED = "#657181"
GRID = "#D9E2EF"
PANEL_FILL = "#FFFFFF"
LIGHT_BLUE = "#F6FAFF"
LIGHT_GREEN = "#F5FBF3"
LIGHT_RED = "#FFF7F5"
GRAY = "#D8DDE5"


def ensure_out_dir():
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def svg_text(x, y, text, size=18, weight="400", fill=INK, anchor="middle", family="Arial", italic=False):
    style = f"font-family:{family}, sans-serif;font-size:{size}px;font-weight:{weight};fill:{fill};"
    if italic:
        style += "font-style:italic;"
    return f'<text x="{x}" y="{y}" text-anchor="{anchor}" style="{style}">{escape(text)}</text>'


def svg_multiline(x, y, lines, size=17, line_gap=22, weight="400", fill=INK, anchor="middle"):
    out = []
    for i, line in enumerate(lines):
        out.append(svg_text(x, y + i * line_gap, line, size=size, weight=weight, fill=fill, anchor=anchor))
    return "\n".join(out)


def svg_rect(x, y, w, h, rx=0, fill="none", stroke=BLUE, sw=2, dash=None, opacity=1):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    return (
        f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="{rx}" '
        f'fill="{fill}" fill-opacity="{opacity}" stroke="{stroke}" stroke-width="{sw}"{dash_attr}/>'
    )


def svg_line(x1, y1, x2, y2, stroke=INK, sw=2, dash=None, marker=False):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    marker_attr = ' marker-end="url(#arrow)"' if marker else ""
    return f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="{stroke}" stroke-width="{sw}"{dash_attr}{marker_attr}/>'


def svg_circle(cx, cy, r, fill=BLUE, stroke="none", sw=1):
    return f'<circle cx="{cx}" cy="{cy}" r="{r}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}"/>'


def svg_path(d, fill="none", stroke=INK, sw=2, dash=None, opacity=1):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    return f'<path d="{d}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}" stroke-opacity="{opacity}" fill-opacity="{opacity}"{dash_attr}/>'


def svg_poly(points, fill="none", stroke=INK, sw=2, opacity=1):
    pts = " ".join(f"{x},{y}" for x, y in points)
    return f'<polygon points="{pts}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}" fill-opacity="{opacity}"/>'


def svg_polyline(points, stroke=INK, sw=2, fill="none", dash=None):
    pts = " ".join(f"{x},{y}" for x, y in points)
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    return f'<polyline points="{pts}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}"{dash_attr}/>'


def add_stage_header(parts, x, y, number, title_lines, w):
    parts.append(svg_circle(x + 20, y + 6, 20, BLUE))
    parts.append(svg_text(x + 20, y + 14, str(number), size=25, weight="700", fill="#FFFFFF"))
    parts.append(svg_multiline(x + w / 2, y + 42, title_lines, size=20, line_gap=24, weight="700"))


def add_arrow(parts, x1, x2, y):
    parts.append(svg_line(x1, y, x2, y, stroke=INK, sw=5, marker=True))


def add_mountain(parts, x, y, w, h, fill="#E8ECEF"):
    pts = [
        (x, y + h),
        (x + w * 0.15, y + h * 0.58),
        (x + w * 0.28, y + h * 0.75),
        (x + w * 0.43, y + h * 0.25),
        (x + w * 0.62, y + h * 0.72),
        (x + w * 0.78, y + h * 0.42),
        (x + w, y + h),
    ]
    parts.append(svg_poly(pts, fill=fill, stroke="#6E7A85", sw=1.5))
    parts.append(svg_polyline([(x + w * 0.16, y + h * 0.58), (x + w * 0.25, y + h * 0.74)], stroke="#A5ADB6", sw=1))
    parts.append(svg_polyline([(x + w * 0.44, y + h * 0.27), (x + w * 0.50, y + h * 0.62)], stroke="#A5ADB6", sw=1))
    parts.append(svg_polyline([(x + w * 0.78, y + h * 0.42), (x + w * 0.70, y + h * 0.73)], stroke="#A5ADB6", sw=1))


def add_grid_risk(parts, x, y, color, label):
    base = [(x, y + 35), (x + 78, y), (x + 156, y + 35), (x + 78, y + 70)]
    parts.append(svg_poly(base, fill="#EEF4F8", stroke="#6E7A85", sw=1))
    for i in range(1, 5):
        parts.append(svg_line(x + i * 18, y + 35 - i * 8, x + 78 + i * 18, y + i * 8, stroke="#B7C3CE", sw=0.8))
        parts.append(svg_line(x + i * 18, y + 35 + i * 8, x + 78 + i * 18, y + 70 - i * 8, stroke="#B7C3CE", sw=0.8))
    hot = [(x + 55, y + 31), (x + 78, y + 20), (x + 101, y + 31), (x + 78, y + 43)]
    parts.append(svg_poly(hot, fill=color, stroke="none", opacity=0.75))
    parts.append(svg_multiline(x + 118, y + 92, label, size=15, line_gap=18, anchor="middle"))


def add_network_plane(parts, cx, cy, w, h, color, label, label_x, label_y):
    pts = [(cx - w / 2, cy), (cx, cy - h / 2), (cx + w / 2, cy), (cx, cy + h / 2)]
    parts.append(svg_poly(pts, fill=color, stroke=color, sw=2, opacity=0.12))
    nodes = [
        (cx - w * 0.30, cy + 1),
        (cx - w * 0.12, cy - h * 0.17),
        (cx + w * 0.03, cy + h * 0.04),
        (cx + w * 0.22, cy - h * 0.10),
        (cx + w * 0.34, cy + h * 0.09),
    ]
    edges = [(0, 1), (1, 2), (2, 3), (3, 4), (0, 2), (2, 4)]
    for a, b in edges:
        parts.append(svg_line(nodes[a][0], nodes[a][1], nodes[b][0], nodes[b][1], stroke=color, sw=2))
    for nx, ny in nodes:
        parts.append(svg_circle(nx, ny, 6, color, stroke="#FFFFFF", sw=1))
    parts.append(svg_multiline(label_x, label_y, label, size=14, line_gap=16, fill=color, anchor="middle"))
    return nodes


def build_svg():
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">',
        "<defs>",
        '<marker id="arrow" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="8" markerHeight="8" orient="auto-start-reverse">',
        f'<path d="M 0 0 L 10 5 L 0 10 z" fill="{INK}"/>',
        "</marker>",
        "</defs>",
        svg_rect(0, 0, W, H, fill="#FFFFFF", stroke="none"),
    ]

    cards = [
        (15, 104, 185, 555),
        (236, 104, 191, 555),
        (452, 104, 205, 555),
        (674, 104, 574, 555),
        (1273, 104, 194, 555),
        (1490, 104, 168, 555),
    ]
    titles = [
        ["Input Data"],
        ["Planning-level", "Risk Fields"],
        ["DEM-driven", "Adaptive Flight", "Corridor"],
        ["Terrain-aware Three-layer Airway Network"],
        ["Event-driven", "Incremental Update"],
        ["Path", "Post-processing"],
    ]

    for i, (x, y, w, h) in enumerate(cards):
        parts.append(svg_rect(x, y, w, h, rx=12, fill=PANEL_FILL, stroke=BLUE_DARK, sw=2))
        add_stage_header(parts, x, y, i + 1, titles[i], w)

    for x1, x2 in [(200, 236), (427, 452), (657, 674), (1248, 1273), (1467, 1490)]:
        add_arrow(parts, x1, x2, 384)

    # 第 1 步：输入数据
    x, y, w, h = cards[0]
    add_mountain(parts, x + 28, y + 118, 96, 58)
    parts.append(svg_text(x + 145, y + 154, "DEM", size=16, weight="700", anchor="middle"))
    for yy in [185, 300, 430]:
        parts.append(svg_line(x + 15, yy, x + w - 15, yy, stroke=GRID, sw=1.4))
    parts.append(svg_poly([(x + 48, y + 224), (x + 102, y + 198), (x + 152, y + 226), (x + 97, y + 254)], fill="#E6F0E4", stroke="#9BA8B6", sw=1.5))
    parts.append(svg_path(f"M{x+102},{y+204} C{x+122},{y+214} {x+122},{y+235} {x+102},{y+248} C{x+82},{y+235} {x+82},{y+214} {x+102},{y+204}Z", fill="#5AA4D7", stroke=BLUE_DARK, sw=1))
    parts.append(svg_circle(x + 102, y + 223, 5, "#FFFFFF"))
    parts.append(svg_multiline(x + 145, y + 237, ["OSM", "features"], size=15, line_gap=18))
    parts.append(svg_poly([(x + 52, y + 332), (x + 110, y + 305), (x + 158, y + 333), (x + 101, y + 360)], fill="#DFF0FF", stroke="#9BA8B6", sw=1))
    parts.append(svg_line(x + 103, y + 318, x + 103, y + 354, stroke=INK, sw=2))
    parts.append(svg_line(x + 88, y + 352, x + 118, y + 352, stroke=INK, sw=2))
    for r in [14, 24]:
        parts.append(svg_path(f"M{x+103-r},{y+321} Q{x+103},{y+306-r/2} {x+103+r},{y+321}", stroke=INK, sw=1.3))
    parts.append(svg_multiline(x + 143, y + 342, ["Communication", "source (LOS)"], size=13.5, line_gap=17))
    parts.append(svg_path(f"M{x+40},{y+448} C{x+60},{y+430} {x+78},{y+449} {x+55},{y+471}", fill="#3F5367", stroke="none"))
    parts.append(svg_circle(x + 55, y + 450, 8, "#FFFFFF"))
    parts.append(svg_path(f"M{x+90},{y+462} C{x+110},{y+438} {x+132},{y+438} {x+152},{y+462}", stroke=INK, sw=1.8, dash="5 5"))
    parts.append(svg_circle(x + 138, y + 446, 9, INK))
    parts.append(svg_line(x + 105, y + 446, x + 170, y + 446, stroke=INK, sw=2))
    parts.append(svg_multiline(x + 142, y + 487, ["Terminals"], size=15))
    parts.append(svg_rect(x + 36, y + 462 + 65, 48, 58, rx=4, fill="#FFFFFF", stroke=INK, sw=2))
    for i in range(3):
        parts.append(svg_line(x + 48, y + 541 + i * 14, x + 72, y + 541 + i * 14, stroke=INK, sw=1.8))
        parts.append(svg_text(x + 42, y + 546 + i * 14, "✓", size=12, fill=BLUE_DARK))
    parts.append(svg_multiline(x + 136, y + 561, ["UAV", "parameters"], size=15, line_gap=18))

    # 第 2 步：风险场
    x, y, w, h = cards[1]
    add_grid_risk(parts, x + 20, y + 120, "#F0C442", ["Terrain", "representation"])
    add_grid_risk(parts, x + 20, y + 270, "#E15845", ["Human exposure", "risk proxy"])
    add_grid_risk(parts, x + 20, y + 420, "#5AA4D7", ["LOS communication", "risk proxy"])
    parts.append(svg_line(x + 92, y + 422, x + 92, y + 458, stroke=INK, sw=2))
    for r in [14, 24]:
        parts.append(svg_path(f"M{x+92-r},{y+427} Q{x+92},{y+410-r/2} {x+92+r},{y+427}", stroke=INK, sw=1.2))

    # 第 3 步：走廊
    x, y, w, h = cards[2]
    add_mountain(parts, x + 35, y + 130, 135, 54)
    parts.append(svg_path(f"M{x+37},{y+139} C{x+70},{y+112} {x+126},{y+126} {x+168},{y+106}", stroke=RED, sw=2, dash="7 4"))
    parts.append(svg_text(x + 93, y + 204, "Floor boundary", size=14, anchor="middle"))
    add_mountain(parts, x + 35, y + 256, 135, 54)
    parts.append(svg_path(f"M{x+37},{y+258} C{x+70},{y+230} {x+126},{y+244} {x+168},{y+224}", stroke=BLUE, sw=2, dash="7 4"))
    parts.append(svg_text(x + 97, y + 334, "Ceiling boundary", size=14, anchor="middle"))
    parts.append(svg_multiline(x + 102, y + 378, ["Three mid-surfaces"], size=16, weight="700"))
    for j, (col, label) in enumerate([(BLUE, "Ceiling mid-surface"), (GREEN, "Middle mid-surface"), (RED, "Floor mid-surface")]):
        yy = y + 410 + j * 30
        parts.append(svg_path(f"M{x+30},{yy} C{x+60},{yy-18} {x+92},{yy+17} {x+124},{yy} S{x+165},{yy-8} {x+182},{yy}", stroke=col, sw=1.8, dash="6 4"))
        parts.append(svg_text(x + 129, yy + 6, label, size=12.5, anchor="start"))
    parts.append(svg_multiline(x + 102, y + 535, ["Flyable mask (corridor)"], size=15, weight="700"))
    parts.append(svg_rect(x + 44, y + 550, 118, 55, rx=0, fill="#F0F2F4", stroke="#7E8790", sw=1.2))
    for gx in range(5):
        parts.append(svg_line(x + 44 + gx * 24, y + 550, x + 44 + gx * 24, y + 605, stroke="#A9B1BA", sw=0.7))
    for gy in range(3):
        parts.append(svg_line(x + 44, y + 550 + gy * 18, x + 162, y + 550 + gy * 18, stroke="#A9B1BA", sw=0.7))
    add_mountain(parts, x + 52, y + 575, 78, 35, fill="#FFFFFF")

    # 第 4 步：三层航路网络核心
    x, y, w, h = cards[3]
    parts.append(svg_text(x + w / 2, y + 36, "Terrain-aware Three-layer Airway Network", size=21, weight="700"))
    add_mountain(parts, x + 70, y + 360, 250, 115)
    n1 = add_network_plane(parts, x + 210, y + 225, 260, 95, BLUE, ["Endpoint access", "layer"], x + 485, y + 235)
    n2 = add_network_plane(parts, x + 210, y + 325, 260, 95, GREEN, ["Regional branch", "layer"], x + 485, y + 340)
    n3 = add_network_plane(parts, x + 210, y + 425, 260, 95, RED, ["Backbone", "layer"], x + 485, y + 452)
    for a, b, c in zip(n1[:4], n2[:4], n3[:4]):
        parts.append(svg_line(a[0], a[1], b[0], b[1], stroke=INK, sw=1.4, dash="6 4"))
        parts.append(svg_line(b[0], b[1], c[0], c[1], stroke=INK, sw=1.4, dash="6 4"))
    parts.append(svg_circle(x + 205, y + 170, 7, INK))
    parts.append(svg_line(x + 185, y + 170, x + 225, y + 170, stroke=INK, sw=2))
    parts.append(svg_line(x + 205, y + 157, x + 205, y + 182, stroke=INK, sw=2))
    parts.append(svg_circle(x + 330, y + 160, 5, BLUE, stroke=INK, sw=1))
    parts.append(svg_line(x + 330, y + 160, x + 330, y + 190, stroke=INK, sw=1.5))
    parts.append(svg_text(x + 45, y + 525, "● Endpoint node", size=12, fill=BLUE, anchor="start"))
    parts.append(svg_text(x + 160, y + 525, "● Branch node", size=12, fill=GREEN, anchor="start"))
    parts.append(svg_text(x + 268, y + 525, "● Backbone node", size=12, fill=RED, anchor="start"))
    parts.append(svg_line(x + 100, y + 545, x + 135, y + 545, stroke=INK, sw=1.3, dash="6 4"))
    parts.append(svg_text(x + 143, y + 550, "Inter-layer link", size=12, anchor="start"))

    mech_x, mech_y, mech_w, mech_h = x + 340, y + 64, 215, 480
    parts.append(svg_rect(mech_x, mech_y, mech_w, mech_h, rx=8, fill="#FFFFFF", stroke=BLUE_DARK, sw=1.8))
    parts.append(svg_text(mech_x + mech_w / 2, mech_y + 22, "Corridor to graph mechanism", size=14.5, weight="700"))
    mech_steps = [
        ("Mid-surfaces input", "wavy"),
        ("Terrain-driven\nnode sampling", "nodes"),
        ("Candidate edges\n(intra-layer)", "edges"),
        ("Safety check\n(clearance & risk)", "shield"),
        ("Inter-layer links\n(vertical connections)", "links"),
        ("Cost assignment\nc(e)", "cost"),
        ("Graph output\nG = (V, E)", "graph"),
    ]
    box_h = 55
    for i, (label, kind) in enumerate(mech_steps):
        by = mech_y + 38 + i * 62
        parts.append(svg_rect(mech_x + 10, by, mech_w - 20, box_h, rx=6, fill=LIGHT_BLUE, stroke="#9CBDE4", sw=1))
        parts.append(svg_circle(mech_x + 23, by + 18, 10, BLUE))
        parts.append(svg_text(mech_x + 23, by + 23, str(i + 1), size=12, weight="700", fill="#FFFFFF"))
        parts.append(svg_multiline(mech_x + 78, by + 22, label.split("\n"), size=12.5, line_gap=15, anchor="middle"))
        ix = mech_x + 145
        iy = by + 18
        if kind == "wavy":
            for off, col in [(0, BLUE), (12, GREEN), (24, RED)]:
                parts.append(svg_path(f"M{ix},{iy+off} C{ix+18},{iy+off-10} {ix+38},{iy+off+10} {ix+58},{iy+off}", stroke=col, sw=1.2, dash="5 3"))
        elif kind == "nodes":
            for j, col in enumerate([BLUE, GREEN, RED]):
                yy = iy + j * 12
                parts.append(svg_line(ix, yy, ix + 60, yy + 8, stroke=col, sw=1.1))
                for k in range(3):
                    parts.append(svg_circle(ix + 8 + k * 22, yy + k * 3, 3, col))
        elif kind == "edges":
            pts = [(ix, iy + 20), (ix + 20, iy + 5), (ix + 40, iy + 23), (ix + 62, iy + 10)]
            parts.append(svg_polyline(pts, stroke=BLUE, sw=1.4))
            for p in pts:
                parts.append(svg_circle(p[0], p[1], 3.5, BLUE))
        elif kind == "shield":
            parts.append(svg_path(f"M{ix+28},{iy} L{ix+55},{iy+10} C{ix+52},{iy+35} {ix+42},{iy+43} {ix+28},{iy+50} C{ix+14},{iy+43} {ix+4},{iy+35} {ix+1},{iy+10}Z", fill="#EAF4FF", stroke=BLUE, sw=2))
            parts.append(svg_text(ix + 28, iy + 31, "✓", size=24, weight="700", fill=BLUE))
        elif kind == "links":
            for k in range(3):
                parts.append(svg_circle(ix + 10 + k * 22, iy + 5, 3, BLUE))
                parts.append(svg_circle(ix + 10 + k * 22, iy + 35, 3, RED))
                parts.append(svg_line(ix + 10 + k * 22, iy + 5, ix + 10 + k * 22, iy + 35, stroke=INK, sw=1.1))
        elif kind == "cost":
            parts.append(svg_rect(ix + 5, iy + 8, 55, 28, rx=4, fill="#FFFFFF", stroke=BLUE_DARK, sw=1.3))
            parts.append(svg_text(ix + 33, iy + 27, "c(e)", size=14, weight="700"))
        else:
            pts = [(ix + 5, iy + 35), (ix + 25, iy + 10), (ix + 45, iy + 32), (ix + 65, iy + 15)]
            parts.append(svg_polyline(pts, stroke=BLUE, sw=1.4))
            for p in pts:
                parts.append(svg_circle(p[0], p[1], 4, BLUE))
        if i < len(mech_steps) - 1:
            parts.append(svg_line(mech_x + mech_w / 2, by + box_h, mech_x + mech_w / 2, by + 62, stroke=INK, sw=1.1, marker=True))

    # 第 5 步：事件驱动增量更新
    x, y, w, h = cards[4]
    parts.append(svg_path(f"M{x+47},{y+148} C{x+35},{y+126} {x+60},{y+108} {x+75},{y+126} C{x+95},{y+116} {x+113},{y+135} {x+98},{y+154} C{x+94},{y+172} {x+57},{y+174} {x+47},{y+148}Z", fill="#B8C6D4", stroke="#526274", sw=2))
    parts.append(svg_text(x + 70, y + 184, "⚡", size=35, fill=ORANGE))
    parts.append(svg_multiline(x + 133, y + 165, ["Regional", "event"], size=14.5, line_gap=18))
    for yy in [290, 430, 558]:
        parts.append(svg_line(x + 12, yy, x + w - 12, yy, stroke=GRID, sw=1.3))
    nodes = [(x + 48, y + 238), (x + 92, y + 222), (x + 118, y + 267), (x + 70, y + 296), (x + 42, y + 275)]
    for a, b in [(0, 1), (1, 2), (2, 3), (3, 4), (4, 0), (0, 3), (1, 3)]:
        col = RED if (a, b) in [(1, 2), (2, 3)] else "#5F6872"
        parts.append(svg_line(nodes[a][0], nodes[a][1], nodes[b][0], nodes[b][1], stroke=col, sw=2 if col == RED else 1.5))
    for i, p in enumerate(nodes):
        parts.append(svg_circle(p[0], p[1], 6, RED if i in [2, 3] else "#5F6872"))
    parts.append(svg_multiline(x + 140, y + 258, ["Local affected", "edges"], size=14, line_gap=18, fill=RED))
    parts.append(svg_rect(x + 47, y + 353, 48, 46, rx=4, fill="#EAF0F4", stroke="#44576A", sw=2))
    for k in range(3):
        parts.append(svg_path(f"M{x+47},{y+360+k*12} C{x+60},{y+368+k*12} {x+82},{y+368+k*12} {x+95},{y+360+k*12}", stroke="#44576A", sw=1.2))
    parts.append(svg_multiline(x + 139, y + 377, ["LPA* state reuse", "(g, rhs)"], size=13.5, line_gap=17))
    parts.append(svg_path(f"M{x+45},{y+514} C{x+78},{y+480} {x+105},{y+540} {x+142},{y+500}", stroke=INK, sw=1.8, dash="5 5"))
    parts.append(svg_circle(x + 45, y + 514, 8, BLUE))
    parts.append(svg_circle(x + 142, y + 500, 8, GREEN))
    parts.append(svg_multiline(x + 133, y + 598, ["Updated path"], size=14))

    # 第 6 步：路径后处理
    x, y, w, h = cards[5]
    for yy in [290, 430, 558]:
        parts.append(svg_line(x + 12, yy, x + w - 12, yy, stroke=GRID, sw=1.3))
    pts = [(x + 38, y + 215), (x + 76, y + 250), (x + 108, y + 218), (x + 142, y + 250)]
    parts.append(svg_polyline(pts, stroke="#5F6872", sw=1.8, dash="6 5"))
    for p in pts:
        parts.append(svg_circle(p[0], p[1], 6, "#5F6872"))
    parts.append(svg_text(x + 92, y + 242, "×", size=24, weight="700", fill=RED))
    parts.append(svg_multiline(x + 88, y + 266, ["LOS pruning"], size=14))
    pts = [(x + 30, y + 380), (x + 70, y + 370), (x + 104, y + 388), (x + 146, y + 356)]
    parts.append(svg_path(f"M{pts[0][0]},{pts[0][1]} C{x+60},{y+350} {x+88},{y+410} {pts[3][0]},{pts[3][1]}", stroke=BLUE, sw=3))
    for p in pts:
        parts.append(svg_circle(p[0], p[1], 6, BLUE))
    parts.append(svg_multiline(x + 88, y + 415, ["B-spline", "smoothing"], size=13.5, line_gap=17))
    parts.append(svg_path(f"M{x+46},{y+468} L{x+75},{y+478} C{x+73},{y+507} {x+60},{y+518} {x+46},{y+526} C{x+32},{y+518} {x+19},{y+507} {x+17},{y+478}Z", fill=LIGHT_GREEN, stroke=GREEN, sw=2.2))
    parts.append(svg_text(x + 46, y + 503, "✓", size=24, weight="700", fill=GREEN))
    parts.append(svg_multiline(x + 116, y + 495, ["Safety recheck", "(clearance & risk)"], size=12.5, line_gap=16))
    parts.append(svg_path(f"M{x+34},{y+610} C{x+75},{y+560} {x+110},{y+650} {x+148},{y+590}", stroke=BLUE, sw=2.8))
    parts.append(svg_circle(x + 34, y + 610, 7, BLUE))
    parts.append(svg_circle(x + 148, y + 590, 7, GREEN))
    parts.append(svg_multiline(x + 92, y + 630, ["Continuous", "trajectory after", "path update"], size=12.5, line_gap=15))

    # 底部贡献说明
    callouts = [
        (207, 695, 409, 80, "Adaptive corridor generation", "mountain"),
        (710, 695, 452, 80, "Structured graph compression", "graph"),
        (1245, 695, 391, 80, "Local update plus post-processing", "update"),
    ]
    targets = [(558, 659), (1018, 659), (1374, 659)]
    for (cx, cy, cw, ch, text, kind), target in zip(callouts, targets):
        parts.append(svg_rect(cx, cy, cw, ch, rx=9, fill="#FFFFFF", stroke=BLUE_DARK, sw=2, dash="8 5"))
        parts.append(svg_line(cx + cw / 2, cy, target[0], target[1], stroke=BLUE_DARK, sw=2, dash="8 5", marker=True))
        if kind == "mountain":
            add_mountain(parts, cx + 48, cy + 20, 85, 42, fill="#EAF4FF")
            parts.append(svg_text(cx + 205, cy + 50, text, size=17, anchor="middle"))
        elif kind == "graph":
            pts = [(cx + 72, cy + 50), (cx + 96, cy + 28), (cx + 124, cy + 52), (cx + 150, cy + 28)]
            parts.append(svg_polyline(pts, stroke=BLUE, sw=2))
            for p in pts:
                parts.append(svg_circle(p[0], p[1], 8, BLUE))
            parts.append(svg_text(cx + 275, cy + 50, text, size=17, anchor="middle"))
        else:
            parts.append(svg_path(f"M{cx+40},{cy+50} C{cx+70},{cy+20} {cx+106},{cy+80} {cx+138},{cy+42}", stroke=INK, sw=1.7, dash="5 5"))
            parts.append(svg_text(cx + 245, cy + 50, text, size=17, anchor="middle"))

    # Outcome
    parts.append(svg_rect(91, 792, 1460, 86, rx=11, fill="#FFFFFF", stroke=BLUE_DARK, sw=2))
    parts.append(svg_circle(160, 835, 43, fill="none", stroke=BLUE, sw=8))
    parts.append(svg_circle(160, 835, 25, fill="none", stroke=BLUE, sw=6))
    parts.append(svg_circle(160, 835, 8, fill=BLUE))
    parts.append(svg_line(160, 835, 201, 794, stroke=BLUE, sw=7, marker=True))
    parts.append(svg_line(330, 806, 330, 864, stroke=BLUE_DARK, sw=2))
    parts.append(svg_text(276, 847, "Outcome", size=22, weight="700", anchor="middle"))
    parts.append(svg_text(887, 846, "Mountainous UAV replanning is transformed into a compact, safety-constrained, incrementally updateable graph search problem.", size=20, anchor="middle"))

    parts.append("</svg>")
    SVG_PATH.write_text("\n".join(parts), encoding="utf-8")


def ppt_color(hex_color):
    return hex_color.replace("#", "")


class PptCanvas:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.sx = 13.333 / W
        self.sy = 7.5 / H

    def ix(self, x):
        return Inches(x * self.sx)

    def iy(self, y):
        return Inches(y * self.sy)

    def rect(self, x, y, w, h, fill="#FFFFFF", stroke=BLUE_DARK, sw=1.2, radius=True, dash=False):
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shp = self.slide.shapes.add_shape(shape_type, self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = pptx_rgb(fill)
        shp.line.color.rgb = pptx_rgb(stroke)
        shp.line.width = Pt(sw)
        if dash:
            shp.line.dash_style = 4
        return shp

    def oval(self, x, y, w, h, fill=BLUE, stroke=None, sw=1):
        shp = self.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = pptx_rgb(fill)
        if stroke:
            shp.line.color.rgb = pptx_rgb(stroke)
            shp.line.width = Pt(sw)
        else:
            shp.line.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
            shp.line.transparency = 100000
        return shp

    def text(self, x, y, w, h, text, size=12, weight=False, color=INK, align=PP_ALIGN.CENTER):
        box = self.slide.shapes.add_textbox(self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        box.text_frame.clear()
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = box.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = weight
        run.font.color.rgb = pptx_rgb(color)
        return box

    def line(self, x1, y1, x2, y2, color=INK, sw=1.2, arrow=False, dash=False):
        conn = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, self.ix(x1), self.iy(y1), self.ix(x2), self.iy(y2))
        conn.line.color.rgb = pptx_rgb(color)
        conn.line.width = Pt(sw)
        if arrow:
            try:
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                from pptx.enum.dml import MSO_ARROWHEAD
                conn.line.end_arrowhead = MSO_ARROWHEAD.TRIANGLE
            except Exception:
                pass
        if dash:
            try:
                from pptx.enum.dml import MSO_LINE_DASH_STYLE
                conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            except Exception:
                pass
        return conn


def pptx_rgb(hex_color):
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def build_pptx():
    # PPT 版本保留主要结构、文本、框线和关键节点，复杂小图标用原生形状简化重建，便于后续人工编辑。
    c = PptCanvas()
    c.rect(0, 0, W, H, fill="#FFFFFF", stroke="#FFFFFF", radius=False)

    cards = [
        (15, 104, 185, 555),
        (236, 104, 191, 555),
        (452, 104, 205, 555),
        (674, 104, 574, 555),
        (1273, 104, 194, 555),
        (1490, 104, 168, 555),
    ]
    titles = [
        "Input Data",
        "Planning-level\nRisk Fields",
        "DEM-driven\nAdaptive Flight\nCorridor",
        "Terrain-aware Three-layer Airway Network",
        "Event-driven\nIncremental Update",
        "Path\nPost-processing",
    ]
    for i, (x, y, w, h) in enumerate(cards):
        c.rect(x, y, w, h, fill="#FFFFFF", stroke=BLUE_DARK, sw=1.2)
        c.oval(x, y - 18, 42, 42, fill=BLUE)
        c.text(x + 4, y - 15, 34, 35, str(i + 1), size=15, weight=True, color="#FFFFFF")
        c.text(x + 10, y + 18, w - 20, 70, titles[i], size=12.2, weight=True)

    for x1, x2 in [(200, 236), (427, 452), (657, 674), (1248, 1273), (1467, 1490)]:
        c.line(x1, 384, x2, 384, color=INK, sw=2.4, arrow=True)

    # 第 1 步到第 3 步的可编辑占位图标和标签
    x, y, w, h = cards[0]
    labels = ["DEM", "OSM\nfeatures", "Communication\nsource (LOS)", "Terminals", "UAV\nparameters"]
    yy = [150, 260, 375, 500, 585]
    for idx, lab in enumerate(labels):
        c.rect(x + 30, yy[idx] - 28, 62, 48, fill="#EEF4F8", stroke="#9BA8B6", sw=0.8, radius=False)
        c.text(x + 100, yy[idx] - 22, 75, 42, lab, size=8.8)
        if idx < 4:
            c.line(x + 15, yy[idx] + 35, x + w - 15, yy[idx] + 35, color=GRID, sw=0.8)

    x, y, w, h = cards[1]
    labels = ["Terrain\nrepresentation", "Human exposure\nrisk proxy", "LOS communication\nrisk proxy"]
    colors = ["#F0C442", "#E15845", "#5AA4D7"]
    for i, lab in enumerate(labels):
        yy0 = y + 135 + i * 150
        c.rect(x + 35, yy0, 85, 55, fill="#EEF4F8", stroke="#9BA8B6", sw=0.8, radius=False)
        c.rect(x + 62, yy0 + 15, 30, 22, fill=colors[i], stroke=colors[i], sw=0.6, radius=False)
        c.text(x + 118, yy0 + 8, 64, 52, lab, size=8.5)

    x, y, w, h = cards[2]
    labels = ["Floor boundary", "Ceiling boundary", "Three mid-surfaces", "Flyable mask\n(corridor)"]
    for i, lab in enumerate(labels):
        yy0 = y + 135 + i * 118
        c.rect(x + 35, yy0, 130, 50, fill="#F6F8FA", stroke="#9BA8B6", sw=0.8, radius=False)
        c.text(x + 30, yy0 + 55, 145, 30, lab, size=8.5)
    c.line(x + 45, y + 414, x + 150, y + 404, color=BLUE, sw=1.0, dash=True)
    c.line(x + 45, y + 444, x + 150, y + 434, color=GREEN, sw=1.0, dash=True)
    c.line(x + 45, y + 474, x + 150, y + 464, color=RED, sw=1.0, dash=True)

    # 第 4 步核心网络，用原生线、圆和文本重建为可编辑对象。
    x, y, w, h = cards[3]
    plane_y = [225, 325, 425]
    plane_colors = [BLUE, GREEN, RED]
    plane_labels = ["Endpoint access\nlayer", "Regional branch\nlayer", "Backbone\nlayer"]
    all_nodes = []
    for py, col, lab in zip(plane_y, plane_colors, plane_labels):
        nodes = [
            (x + 130, py),
            (x + 190, py - 33),
            (x + 235, py + 5),
            (x + 300, py - 22),
            (x + 365, py + 10),
        ]
        all_nodes.append(nodes)
        for a, b in [(0, 1), (1, 2), (2, 3), (3, 4), (0, 2), (2, 4)]:
            c.line(nodes[a][0], nodes[a][1], nodes[b][0], nodes[b][1], color=col, sw=1.1)
        for nx, ny in nodes:
            c.oval(nx - 6, ny - 6, 12, 12, fill=col, stroke="#FFFFFF", sw=0.5)
        c.text(x + 430, py - 20, 105, 42, lab, size=8, color=col)
    for layer in range(2):
        for i in range(4):
            c.line(all_nodes[layer][i][0], all_nodes[layer][i][1], all_nodes[layer + 1][i][0], all_nodes[layer + 1][i][1], color=INK, sw=0.7, dash=True)
    c.rect(x + 338, y + 64, 218, 480, fill="#FFFFFF", stroke=BLUE_DARK, sw=1.0)
    c.text(x + 350, y + 66, 194, 28, "Corridor to graph mechanism", size=8.2, weight=True)
    mech = [
        "1  Mid-surfaces input",
        "2  Terrain-driven\n    node sampling",
        "3  Candidate edges\n    intra-layer",
        "4  Safety check\n    clearance & risk",
        "5  Inter-layer links",
        "6  Cost assignment\n    c(e)",
        "7  Graph output\n    G=(V,E)",
    ]
    for i, lab in enumerate(mech):
        by = y + 105 + i * 61
        c.rect(x + 350, by, 190, 43, fill=LIGHT_BLUE, stroke="#9CBDE4", sw=0.6)
        c.text(x + 360, by + 3, 165, 36, lab, size=6.8, align=PP_ALIGN.LEFT)
        if i < 6:
            c.line(x + 445, by + 43, x + 445, by + 58, color=INK, sw=0.55, arrow=True)

    # 第 5 步、第 6 步的可编辑模块
    x, y, w, h = cards[4]
    c.rect(x + 40, y + 120, 65, 58, fill="#DDE6EE", stroke="#526274", sw=1.0)
    c.text(x + 118, y + 138, 58, 48, "Regional\nevent", size=8.5)
    c.line(x + 15, y + 290, x + w - 15, y + 290, color=GRID, sw=0.8)
    c.line(x + 15, y + 430, x + w - 15, y + 430, color=GRID, sw=0.8)
    c.line(x + 15, y + 558, x + w - 15, y + 558, color=GRID, sw=0.8)
    gn = [(x + 48, y + 238), (x + 92, y + 222), (x + 118, y + 267), (x + 70, y + 296), (x + 42, y + 275)]
    for a, b in [(0, 1), (1, 2), (2, 3), (3, 4), (4, 0), (0, 3)]:
        c.line(gn[a][0], gn[a][1], gn[b][0], gn[b][1], color=RED if (a, b) in [(1, 2), (2, 3)] else "#5F6872", sw=1.0)
    for i, p in enumerate(gn):
        c.oval(p[0] - 5, p[1] - 5, 10, 10, fill=RED if i in [2, 3] else "#5F6872")
    c.text(x + 126, y + 244, 55, 45, "Local affected\nedges", size=7.5, color=RED)
    c.rect(x + 48, y + 352, 48, 46, fill="#EAF0F4", stroke="#44576A", sw=1.0)
    c.text(x + 112, y + 360, 70, 50, "LPA* state reuse\n(g, rhs)", size=7.2)
    c.line(x + 45, y + 514, x + 142, y + 500, color=INK, sw=1.0, dash=True)
    c.oval(x + 37, y + 506, 16, 16, fill=BLUE)
    c.oval(x + 134, y + 492, 16, 16, fill=GREEN)
    c.text(x + 102, y + 578, 80, 36, "Updated path", size=8.2)

    x, y, w, h = cards[5]
    for yy0 in [290, 430, 558]:
        c.line(x + 12, yy0, x + w - 12, yy0, color=GRID, sw=0.8)
    c.line(x + 38, y + 215, x + 142, y + 250, color="#5F6872", sw=1.0, dash=True)
    c.text(x + 35, y + 248, 120, 30, "LOS pruning", size=8.2)
    c.line(x + 30, y + 380, x + 146, y + 356, color=BLUE, sw=1.4)
    c.text(x + 35, y + 400, 120, 45, "B-spline\nsmoothing", size=7.8)
    c.rect(x + 25, y + 468, 48, 58, fill=LIGHT_GREEN, stroke=GREEN, sw=1.1)
    c.text(x + 80, y + 480, 80, 50, "Safety recheck\n(clearance & risk)", size=7)
    c.line(x + 34, y + 610, x + 148, y + 590, color=BLUE, sw=1.4)
    c.text(x + 30, y + 620, 125, 52, "Continuous\ntrajectory after\npath update", size=7.2)

    # 底部 callout 和 outcome
    for cx, cy, cw, ch, text in [
        (207, 695, 409, 80, "Adaptive corridor generation"),
        (710, 695, 452, 80, "Structured graph compression"),
        (1245, 695, 391, 80, "Local update plus post-processing"),
    ]:
        c.rect(cx, cy, cw, ch, fill="#FFFFFF", stroke=BLUE_DARK, sw=1.0, dash=True)
        c.text(cx + 120, cy + 20, cw - 140, 40, text, size=10.5)
    c.rect(91, 792, 1460, 86, fill="#FFFFFF", stroke=BLUE_DARK, sw=1.2)
    c.oval(125, 805, 70, 70, fill="#FFFFFF", stroke=BLUE, sw=3)
    c.text(236, 810, 120, 48, "Outcome", size=14, weight=True)
    c.text(390, 812, 1060, 48, "Mountainous UAV replanning is transformed into a compact, safety-constrained, incrementally updateable graph search problem.", size=12)

    c.prs.save(PPTX_PATH)


def write_manifest():
    MANIFEST_PATH.write_text(
        "\n".join(
            [
                "# 最终图可编辑化产物说明",
                "",
                f"源图：`../S5-candidate-images/candidate-03.png`",
                f"可编辑 SVG：`{SVG_PATH.name}`",
                f"可编辑 PPTX：`{PPTX_PATH.name}`",
                "",
                "## 转换策略",
                "",
                "本次不是把原始 PNG 整页嵌入，而是按最终图的论文语义进行对象级重绘。SVG 中的卡片、箭头、节点、图标、符号和文字均为可编辑矢量或文本对象；PPTX 中主要结构、文本、线条、圆点和机制框为 PowerPoint 原生对象。",
                "",
                "## 可编辑范围",
                "",
                "可直接编辑阶段标题、模块标签、caption 相关术语、三层节点颜色、局部受影响边颜色、箭头粗细、卡片尺寸和底部 outcome 文本。",
                "",
                "## 已保留的关键语义",
                "",
                "第 4 步仍为视觉核心，蓝色、绿色和红色节点面分别表示端点接入层、区域支路层和骨干通行层；第 5 步红色局部边表示区域事件影响的受影响边集合；第 6 步为路径后处理，图内编号为 6，没有写成 3.6。",
                "",
                "## 已知取舍",
                "",
                "SVG 版本更接近最终 raster 图的视觉细节；PPTX 版本优先对象级可编辑性，复杂小图标被简化为原生形状，适合后续在 PowerPoint 中继续微调。",
                "",
            ]
        ),
        encoding="utf-8",
    )


def main():
    ensure_out_dir()
    build_svg()
    build_pptx()
    write_manifest()
    print(SVG_PATH)
    print(PPTX_PATH)
    print(MANIFEST_PATH)


if __name__ == "__main__":
    main()
