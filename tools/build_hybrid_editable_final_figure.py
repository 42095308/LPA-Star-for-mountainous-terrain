from base64 import b64encode
from pathlib import Path
from xml.sax.saxutils import escape

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


W, H = 1672, 941
SRC = Path(
    "intermediate_artifacts/paper_framework_figure_studio/"
    "third_chapter_method_framework/outputs/S5-candidate-images/candidate-03.png"
)
OUT_DIR = Path(
    "intermediate_artifacts/paper_framework_figure_studio/"
    "third_chapter_method_framework/outputs/editable-final-figure-hybrid"
)
ASSET_DIR = OUT_DIR / "assets"
SVG_PATH = OUT_DIR / "hybrid_editable_final_framework.svg"
PPTX_PATH = OUT_DIR / "hybrid_editable_final_framework.pptx"
MANIFEST_PATH = OUT_DIR / "hybrid_editable_conversion_manifest.md"

BLUE = "#1557A8"
BLUE_DARK = "#0B3F86"
INK = "#1D2733"
MUTED = "#657181"
GRID = "#D9E2EF"
WHITE = "#FFFFFF"


# 复杂或难以可靠转换的绘制区域，按原图局部裁切保留。
# x, y, w, h 是源图像素坐标，目标位置默认与源图一致。
CROPS = [
    ("stage1_dem", 34, 205, 105, 70),
    ("stage1_osm", 38, 302, 108, 88),
    ("stage1_los", 45, 410, 108, 78),
    ("stage1_terminals", 28, 522, 124, 66),
    ("stage1_uav", 36, 602, 70, 64),
    ("stage2_terrain", 248, 214, 126, 92),
    ("stage2_exposure", 249, 360, 126, 96),
    ("stage2_los", 249, 504, 126, 98),
    ("stage3_floor", 468, 206, 145, 96),
    ("stage3_ceiling", 469, 315, 145, 96),
    ("stage3_mids", 470, 437, 142, 93),
    ("stage3_mask", 467, 565, 152, 88),
    ("stage4_network", 687, 170, 316, 405),
    ("stage4_mech_1", 1134, 202, 76, 42),
    ("stage4_mech_2", 1134, 280, 76, 40),
    ("stage4_mech_3", 1134, 355, 76, 38),
    ("stage4_mech_4", 1148, 414, 54, 58),
    ("stage4_mech_5", 1136, 487, 76, 48),
    ("stage4_mech_6", 1150, 552, 58, 36),
    ("stage4_mech_7", 1135, 616, 78, 38),
    ("stage5_event", 1306, 244, 66, 70),
    ("stage5_edges", 1282, 334, 92, 82),
    ("stage5_state", 1280, 482, 70, 68),
    ("stage5_path", 1290, 604, 80, 60),
    ("stage6_los", 1514, 204, 118, 72),
    ("stage6_bspline", 1510, 360, 124, 70),
    ("stage6_shield", 1512, 486, 68, 72),
    ("stage6_traj", 1514, 593, 116, 80),
    ("callout_corridor", 232, 716, 112, 54),
    ("callout_graph", 755, 720, 106, 58),
    ("callout_update", 1262, 718, 100, 58),
    ("outcome_target", 118, 806, 105, 76),
]


def ensure_dirs():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def crop_assets():
    im = Image.open(SRC).convert("RGBA")
    assets = {}
    for name, x, y, w, h in CROPS:
        crop = im.crop((x, y, x + w, y + h))
        path = ASSET_DIR / f"{name}.png"
        crop.save(path)
        assets[name] = {"path": path, "x": x, "y": y, "w": w, "h": h}
    return assets


def image_data_uri(path):
    data = b64encode(path.read_bytes()).decode("ascii")
    return f"data:image/png;base64,{data}"


def svg_text(x, y, text, size=18, weight="400", fill=INK, anchor="middle"):
    style = f"font-family:Arial, Helvetica, sans-serif;font-size:{size}px;font-weight:{weight};fill:{fill};"
    return f'<text x="{x}" y="{y}" text-anchor="{anchor}" style="{style}">{escape(text)}</text>'


def svg_multiline(x, y, lines, size=16, gap=20, weight="400", fill=INK, anchor="middle"):
    return "\n".join(
        svg_text(x, y + i * gap, line, size=size, weight=weight, fill=fill, anchor=anchor)
        for i, line in enumerate(lines)
    )


def svg_rect(x, y, w, h, rx=0, fill=WHITE, stroke=BLUE_DARK, sw=2, dash=None):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    return f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="{rx}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}"{dash_attr}/>'


def svg_line(x1, y1, x2, y2, stroke=INK, sw=2, dash=None, marker=False):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    marker_attr = ' marker-end="url(#arrow)"' if marker else ""
    return f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="{stroke}" stroke-width="{sw}"{dash_attr}{marker_attr}/>'


def svg_circle(cx, cy, r, fill=BLUE, stroke="none", sw=1):
    return f'<circle cx="{cx}" cy="{cy}" r="{r}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}"/>'


def svg_image(asset):
    return (
        f'<image x="{asset["x"]}" y="{asset["y"]}" width="{asset["w"]}" height="{asset["h"]}" '
        f'href="{image_data_uri(asset["path"])}"/>'
    )


def add_header(parts, x, y, w, number, title_lines):
    parts.append(svg_circle(x + 20, y + 6, 20, BLUE))
    parts.append(svg_text(x + 20, y + 14, str(number), size=25, weight="700", fill=WHITE))
    parts.append(svg_multiline(x + w / 2, y + 42, title_lines, size=20, gap=24, weight="700"))


def add_arrow(parts, x1, x2, y):
    parts.append(svg_line(x1, y, x2, y, stroke=INK, sw=5, marker=True))


def build_svg(assets):
    cards = [
        (15, 104, 185, 555),
        (236, 104, 191, 555),
        (452, 104, 205, 555),
        (674, 104, 574, 555),
        (1273, 104, 194, 555),
        (1490, 104, 168, 555),
    ]
    titles = [
        ["Input Data"],
        ["Planning-level", "Risk Fields"],
        ["DEM-driven", "Adaptive Flight", "Corridor"],
        ["Terrain-aware Three-layer Airway Network"],
        ["Event-driven", "Incremental Update"],
        ["Path", "Post-processing"],
    ]

    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">',
        "<defs>",
        '<marker id="arrow" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="8" markerHeight="8" orient="auto-start-reverse">',
        f'<path d="M 0 0 L 10 5 L 0 10 z" fill="{INK}"/>',
        "</marker>",
        "</defs>",
        svg_rect(0, 0, W, H, fill=WHITE, stroke=WHITE),
    ]

    for i, (x, y, w, h) in enumerate(cards):
        parts.append(svg_rect(x, y, w, h, rx=12, fill=WHITE, stroke=BLUE_DARK, sw=2))
        add_header(parts, x, y, w, i + 1, titles[i])

    for x1, x2 in [(200, 236), (427, 452), (657, 674), (1248, 1273), (1467, 1490)]:
        add_arrow(parts, x1, x2, 384)

    for name in [
        "stage1_dem", "stage1_osm", "stage1_los", "stage1_terminals", "stage1_uav",
        "stage2_terrain", "stage2_exposure", "stage2_los",
        "stage3_floor", "stage3_ceiling", "stage3_mids", "stage3_mask",
        "stage4_network",
        "stage5_event", "stage5_edges", "stage5_state", "stage5_path",
        "stage6_los", "stage6_bspline", "stage6_shield", "stage6_traj",
        "callout_corridor", "callout_graph", "callout_update", "outcome_target",
    ]:
        parts.append(svg_image(assets[name]))

    # 可编辑分隔线
    for x, ys in [(15, [185, 300, 430]), (1273, [290, 430, 558]), (1490, [290, 430, 558])]:
        w = 185 if x == 15 else 194 if x == 1273 else 168
        for yy in ys:
            parts.append(svg_line(x + 12, yy, x + w - 12, yy, stroke=GRID, sw=1.3))

    # 第 1 步可编辑文字
    x, y, w, h = cards[0]
    parts.append(svg_text(x + 145, y + 154, "DEM", size=16, weight="700"))
    parts.append(svg_multiline(x + 145, y + 237, ["OSM", "features"], size=15, gap=18))
    parts.append(svg_multiline(x + 143, y + 342, ["Communication", "source (LOS)"], size=13.5, gap=17))
    parts.append(svg_text(x + 142, y + 487, "Terminals", size=15))
    parts.append(svg_multiline(x + 136, y + 561, ["UAV", "parameters"], size=15, gap=18))

    # 第 2 步可编辑文字
    x, y, w, h = cards[1]
    parts.append(svg_multiline(x + 150, y + 245, ["Terrain", "representation"], size=15, gap=18))
    parts.append(svg_multiline(x + 150, y + 394, ["Human exposure", "risk proxy"], size=14.5, gap=18))
    parts.append(svg_multiline(x + 150, y + 548, ["LOS communication", "risk proxy"], size=13.5, gap=17))

    # 第 3 步可编辑文字
    x, y, w, h = cards[2]
    parts.append(svg_text(x + 90, y + 204, "Floor boundary", size=14))
    parts.append(svg_text(x + 98, y + 334, "Ceiling boundary", size=14))
    parts.append(svg_text(x + 102, y + 378, "Three mid-surfaces", size=16, weight="700"))
    parts.append(svg_text(x + 132, y + 422, "Ceiling mid-surface", size=12.5, anchor="start"))
    parts.append(svg_text(x + 132, y + 452, "Middle mid-surface", size=12.5, anchor="start"))
    parts.append(svg_text(x + 132, y + 482, "Floor mid-surface", size=12.5, anchor="start"))
    parts.append(svg_text(x + 102, y + 535, "Flyable mask (corridor)", size=15, weight="700"))

    # 第 4 步机制面板，文字与线框可编辑，小图按原图贴片保留
    x, y, w, h = cards[3]
    mech_x, mech_y, mech_w, mech_h = x + 340, y + 64, 215, 480
    parts.append(svg_rect(mech_x, mech_y, mech_w, mech_h, rx=8, fill=WHITE, stroke=BLUE_DARK, sw=1.8))
    parts.append(svg_text(mech_x + mech_w / 2, mech_y + 22, "Corridor to graph mechanism", size=14.5, weight="700"))
    mech_steps = [
        ("Mid-surfaces input", "stage4_mech_1"),
        ("Terrain-driven\nnode sampling", "stage4_mech_2"),
        ("Candidate edges\n(intra-layer)", "stage4_mech_3"),
        ("Safety check\n(clearance & risk)", "stage4_mech_4"),
        ("Inter-layer links\n(vertical connections)", "stage4_mech_5"),
        ("Cost assignment\nc(e)", "stage4_mech_6"),
        ("Graph output\nG = (V, E)", "stage4_mech_7"),
    ]
    for i, (label, asset_name) in enumerate(mech_steps):
        by = mech_y + 38 + i * 62
        parts.append(svg_rect(mech_x + 10, by, mech_w - 20, 55, rx=6, fill="#F6FAFF", stroke="#9CBDE4", sw=1))
        parts.append(svg_circle(mech_x + 23, by + 18, 10, BLUE))
        parts.append(svg_text(mech_x + 23, by + 23, str(i + 1), size=12, weight="700", fill=WHITE))
        parts.append(svg_multiline(mech_x + 78, by + 22, label.split("\n"), size=12.5, gap=15))
        parts.append(svg_image(assets[asset_name]))
        if i < len(mech_steps) - 1:
            parts.append(svg_line(mech_x + mech_w / 2, by + 55, mech_x + mech_w / 2, by + 62, stroke=INK, sw=1.1, marker=True))

    # 第 5 步可编辑文字
    x, y, w, h = cards[4]
    parts.append(svg_multiline(x + 133, y + 165, ["Regional", "event"], size=14.5, gap=18))
    parts.append(svg_multiline(x + 140, y + 258, ["Local affected", "edges"], size=14, gap=18, fill="#D42E20"))
    parts.append(svg_multiline(x + 139, y + 377, ["LPA* state reuse", "(g, rhs)"], size=13.5, gap=17))
    parts.append(svg_text(x + 133, y + 598, "Updated path", size=14))

    # 第 6 步可编辑文字
    x, y, w, h = cards[5]
    parts.append(svg_text(x + 88, y + 266, "LOS pruning", size=14))
    parts.append(svg_multiline(x + 88, y + 415, ["B-spline", "smoothing"], size=13.5, gap=17))
    parts.append(svg_multiline(x + 116, y + 495, ["Safety recheck", "(clearance & risk)"], size=12.5, gap=16))
    parts.append(svg_multiline(x + 92, y + 630, ["Continuous", "trajectory after", "path update"], size=12.5, gap=15))

    # 底部 callouts 与 outcome
    callouts = [
        (207, 695, 409, 80, "Adaptive corridor generation", 345),
        (710, 695, 452, 80, "Structured graph compression", 930),
        (1245, 695, 391, 80, "Local update plus post-processing", 1490),
    ]
    for cx, cy, cw, ch, text, tx in callouts:
        parts.append(svg_rect(cx, cy, cw, ch, rx=9, fill=WHITE, stroke=BLUE_DARK, sw=2, dash="8 5"))
        parts.append(svg_text(tx, cy + 50, text, size=17))
    parts.append(svg_line(207 + 409 / 2, 695, 558, 659, stroke=BLUE_DARK, sw=2, dash="8 5", marker=True))
    parts.append(svg_line(710 + 452 / 2, 695, 1018, 659, stroke=BLUE_DARK, sw=2, dash="8 5", marker=True))
    parts.append(svg_line(1245 + 391 / 2, 695, 1374, 659, stroke=BLUE_DARK, sw=2, dash="8 5", marker=True))

    parts.append(svg_rect(91, 792, 1460, 86, rx=11, fill=WHITE, stroke=BLUE_DARK, sw=2))
    parts.append(svg_line(330, 806, 330, 864, stroke=BLUE_DARK, sw=2))
    parts.append(svg_text(276, 847, "Outcome", size=22, weight="700"))
    parts.append(svg_text(887, 846, "Mountainous UAV replanning is transformed into a compact, safety-constrained, incrementally updateable graph search problem.", size=20))
    parts.append("</svg>")
    SVG_PATH.write_text("\n".join(parts), encoding="utf-8")


def pptx_rgb(hex_color):
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Canvas:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.sx = 13.333 / W
        self.sy = 7.5 / H

    def ix(self, x):
        return Inches(x * self.sx)

    def iy(self, y):
        return Inches(y * self.sy)

    def rect(self, x, y, w, h, fill=WHITE, stroke=BLUE_DARK, sw=1.2, radius=True):
        typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shp = self.slide.shapes.add_shape(typ, self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = pptx_rgb(fill)
        shp.line.color.rgb = pptx_rgb(stroke)
        shp.line.width = Pt(sw)
        return shp

    def image(self, asset):
        self.slide.shapes.add_picture(str(asset["path"]), self.ix(asset["x"]), self.iy(asset["y"]), self.ix(asset["w"]), self.iy(asset["h"]))

    def oval(self, x, y, w, h, fill=BLUE):
        shp = self.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = pptx_rgb(fill)
        shp.line.color.rgb = pptx_rgb(fill)
        return shp

    def text(self, x, y, w, h, text, size=10, weight=False, color=INK, align=PP_ALIGN.CENTER):
        box = self.slide.shapes.add_textbox(self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        box.text_frame.clear()
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = box.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = weight
        run.font.color.rgb = pptx_rgb(color)
        return box

    def line(self, x1, y1, x2, y2, color=INK, sw=1.2):
        conn = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, self.ix(x1), self.iy(y1), self.ix(x2), self.iy(y2))
        conn.line.color.rgb = pptx_rgb(color)
        conn.line.width = Pt(sw)
        return conn


def build_pptx(assets):
    c = Canvas()
    c.rect(0, 0, W, H, fill=WHITE, stroke=WHITE, radius=False)
    cards = [
        (15, 104, 185, 555),
        (236, 104, 191, 555),
        (452, 104, 205, 555),
        (674, 104, 574, 555),
        (1273, 104, 194, 555),
        (1490, 104, 168, 555),
    ]
    titles = [
        "Input Data",
        "Planning-level\nRisk Fields",
        "DEM-driven\nAdaptive Flight\nCorridor",
        "Terrain-aware Three-layer Airway Network",
        "Event-driven\nIncremental Update",
        "Path\nPost-processing",
    ]
    for i, (x, y, w, h) in enumerate(cards):
        c.rect(x, y, w, h, fill=WHITE, stroke=BLUE_DARK)
        c.oval(x, y - 18, 42, 42, fill=BLUE)
        c.text(x + 4, y - 15, 34, 35, str(i + 1), size=15, weight=True, color=WHITE)
        c.text(x + 10, y + 18, w - 20, 70, titles[i], size=12, weight=True)

    for x1, x2 in [(200, 236), (427, 452), (657, 674), (1248, 1273), (1467, 1490)]:
        c.line(x1, 384, x2, 384, color=INK, sw=2.4)

    # 先放复杂贴片，再放可编辑文字，避免文字被贴片覆盖。
    for name in [
        "stage1_dem", "stage1_osm", "stage1_los", "stage1_terminals", "stage1_uav",
        "stage2_terrain", "stage2_exposure", "stage2_los",
        "stage3_floor", "stage3_ceiling", "stage3_mids", "stage3_mask",
        "stage4_network",
        "stage5_event", "stage5_edges", "stage5_state", "stage5_path",
        "stage6_los", "stage6_bspline", "stage6_shield", "stage6_traj",
        "callout_corridor", "callout_graph", "callout_update", "outcome_target",
    ]:
        c.image(assets[name])

    for x, ys in [(15, [185, 300, 430]), (1273, [290, 430, 558]), (1490, [290, 430, 558])]:
        w = 185 if x == 15 else 194 if x == 1273 else 168
        for yy in ys:
            c.line(x + 12, yy, x + w - 12, yy, color=GRID, sw=0.7)

    labels = [
        (160, 241, 60, 26, "DEM", 9, True, INK),
        (130, 320, 62, 45, "OSM\nfeatures", 8.5, False, INK),
        (120, 438, 76, 45, "Communication\nsource (LOS)", 7.5, False, INK),
        (117, 570, 75, 30, "Terminals", 8.5, False, INK),
        (106, 646, 82, 45, "UAV\nparameters", 8.5, False, INK),
        (346, 242, 70, 46, "Terrain\nrepresentation", 8.2, False, INK),
        (338, 392, 80, 46, "Human exposure\nrisk proxy", 8.0, False, INK),
        (336, 548, 82, 46, "LOS communication\nrisk proxy", 7.5, False, INK),
        (472, 190, 120, 28, "Floor boundary", 8.2, False, INK),
        (473, 319, 130, 28, "Ceiling boundary", 8.2, False, INK),
        (463, 366, 170, 30, "Three mid-surfaces", 9, True, INK),
        (572, 421, 85, 22, "Ceiling mid-surface", 6.6, False, INK),
        (572, 451, 85, 22, "Middle mid-surface", 6.6, False, INK),
        (572, 481, 85, 22, "Floor mid-surface", 6.6, False, INK),
        (468, 520, 150, 28, "Flyable mask (corridor)", 8.2, True, INK),
        (1368, 150, 70, 45, "Regional\nevent", 8.2, False, INK),
        (1360, 242, 86, 50, "Local affected\nedges", 8.0, False, "#D42E20"),
        (1355, 355, 90, 55, "LPA* state reuse\n(g, rhs)", 7.5, False, INK),
        (1350, 580, 95, 32, "Updated path", 8.0, False, INK),
        (1530, 248, 105, 28, "LOS pruning", 8.0, False, INK),
        (1530, 395, 102, 42, "B-spline\nsmoothing", 7.6, False, INK),
        (1578, 474, 72, 50, "Safety recheck\n(clearance & risk)", 6.5, False, INK),
        (1525, 620, 115, 52, "Continuous\ntrajectory after\npath update", 6.8, False, INK),
    ]
    for x, y, w, h, t, s, b, col in labels:
        c.text(x, y, w, h, t, size=s, weight=b, color=col)

    # 第 4 步机制面板可编辑线框和文字，右侧小图保持原样贴片。
    mech_x, mech_y, mech_w, mech_h = 1014, 168, 215, 480
    c.rect(mech_x, mech_y, mech_w, mech_h, fill=WHITE, stroke=BLUE_DARK)
    c.text(mech_x + 8, mech_y + 4, mech_w - 16, 28, "Corridor to graph mechanism", size=8.2, weight=True)
    mech_steps = [
        ("Mid-surfaces input", "stage4_mech_1"),
        ("Terrain-driven\nnode sampling", "stage4_mech_2"),
        ("Candidate edges\n(intra-layer)", "stage4_mech_3"),
        ("Safety check\n(clearance & risk)", "stage4_mech_4"),
        ("Inter-layer links\n(vertical connections)", "stage4_mech_5"),
        ("Cost assignment\nc(e)", "stage4_mech_6"),
        ("Graph output\nG = (V, E)", "stage4_mech_7"),
    ]
    for i, (txt, aname) in enumerate(mech_steps):
        by = mech_y + 38 + i * 62
        c.rect(mech_x + 10, by, mech_w - 20, 55, fill="#F6FAFF", stroke="#9CBDE4")
        c.oval(mech_x + 13, by + 8, 20, 20, fill=BLUE)
        c.text(mech_x + 14, by + 8, 18, 18, str(i + 1), size=6.5, weight=True, color=WHITE)
        c.text(mech_x + 40, by + 8, 95, 38, txt, size=6.8)
        c.image(assets[aname])

    for cx, cy, cw, ch, text, tx in [
        (207, 695, 409, 80, "Adaptive corridor generation", 345),
        (710, 695, 452, 80, "Structured graph compression", 930),
        (1245, 695, 391, 80, "Local update plus post-processing", 1490),
    ]:
        c.rect(cx, cy, cw, ch, fill=WHITE, stroke=BLUE_DARK)
        c.text(tx - 165, cy + 22, 330, 36, text, size=9.5)

    c.rect(91, 792, 1460, 86, fill=WHITE, stroke=BLUE_DARK)
    c.line(330, 806, 330, 864, color=BLUE_DARK, sw=1)
    c.text(235, 815, 100, 40, "Outcome", size=13, weight=True)
    c.text(390, 812, 1060, 48, "Mountainous UAV replanning is transformed into a compact, safety-constrained, incrementally updateable graph search problem.", size=11)

    c.prs.save(PPTX_PATH)


def write_manifest(assets):
    converted = [
        "阶段编号圆点、六个阶段外框、主流程箭头、底部分组框、Outcome 框",
        "主要英文标题、模块标签、机制链文字、legend 相关术语标签",
        "第 5 步与第 6 步的分隔线、底部贡献说明连接线",
    ]
    preserved = [
        "DEM、OSM、通信源、端点、UAV 参数图标",
        "三类风险面、走廊边界和三层中面示意",
        "第 4 步三层航路网络主体和机制链右侧小图",
        "区域事件、局部边、状态复用、路径后处理图标",
        "底部 callout 图标和 outcome 靶标图标",
    ]
    asset_lines = [f"- `{name}.png`，源图坐标 x={a['x']}，y={a['y']}，w={a['w']}，h={a['h']}" for name, a in assets.items()]
    MANIFEST_PATH.write_text(
        "\n".join(
            [
                "# 混合可编辑化产物说明",
                "",
                f"源图：`{SRC}`",
                f"混合可编辑 SVG：`{SVG_PATH.name}`",
                f"混合可编辑 PPTX：`{PPTX_PATH.name}`",
                "",
                "## 转换原则",
                "",
                "本版本采用混合策略：可以稳定转换的文字、线条、外框、箭头和编号转为可编辑对象；难以准确转换的绘制型区域直接从最终 PNG 中裁切并按原位置放置，保留原样视觉。",
                "",
                "## 已转为可编辑对象",
                "",
                *[f"- {item}" for item in converted],
                "",
                "## 按原样保留的 raster 局部",
                "",
                *[f"- {item}" for item in preserved],
                "",
                "## 局部贴片清单",
                "",
                *asset_lines,
                "",
                "## 语义检查",
                "",
                "第 4 步三层航路网络仍是视觉核心；第 4 步中蓝色、绿色和红色节点面表示端点接入层、区域支路层和骨干通行层。第 5 步红色局部边表示区域事件影响的受影响边集合。第 6 步仍为路径后处理，图内编号为 6。",
                "",
                "## 使用建议",
                "",
                "需要改文字、箭头、框线、编号和模块层级时，优先使用 PPTX 或 SVG 直接编辑。需要改复杂网络主体、风险面或图标时，建议替换对应 `assets` 局部贴片，或回到图像生成步骤局部重绘。",
                "",
            ]
        ),
        encoding="utf-8",
    )


def main():
    ensure_dirs()
    assets = crop_assets()
    build_svg(assets)
    build_pptx(assets)
    write_manifest(assets)
    print(SVG_PATH)
    print(PPTX_PATH)
    print(MANIFEST_PATH)
    print(f"assets={len(assets)}")


if __name__ == "__main__":
    main()
